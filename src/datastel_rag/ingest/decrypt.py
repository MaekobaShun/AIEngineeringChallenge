"""Generic password derivation/decryption for MS-CFB-encrypted Office files.

Rule (from 社内管理/データアステル社内規定_パスワード導出規則.docx):

    DA-[案件略号]-[開始年月日8桁]-[拡張子コード]

- 案件略号: the project's 主略称 from the glossary (catalog.glossary).
- 開始年月日8桁: the project's contract start date, YYYYMMDD. Not every
  project states this in an obviously-labeled field, so we harvest date
  candidates from every *readable* sibling file in the same project folder
  (filenames + text content) and try them all.
- 拡張子コード: assumed to be the file extension, but we don't know the
  exact casing/format the puzzle wants, so several variants are tried.

Never hardcode a specific project's password. This module only encodes the
*algorithm*; candidates are always derived at runtime from the glossary and
sibling documents, and it degrades gracefully (returns None) when nothing
in the candidate space works -- callers (including the agent, via a tool)
can then extend the candidate list with their own guesses.
"""

from __future__ import annotations

import io
import itertools
import json
import re
from dataclasses import dataclass

import docx
import msoffcrypto
from pptx import Presentation

from datastel_rag import config
from datastel_rag.catalog.glossary import Glossary
from datastel_rag.catalog.scanner import Catalog, FileEntry
from datastel_rag.ingest.fileformats import sniff_is_ole

_DATE_RE = re.compile(r"(20\d{2})[年/\-.](\d{1,2})[月/\-.](\d{1,2})")

_PASSWORD_CACHE_PATH = config.CACHE_DIR / "decrypt_passwords.json"


def _extract_plain_text(entry: FileEntry) -> str:
    """Best-effort plain text for date harvesting only -- not the real parser."""
    try:
        if entry.ext == "docx":
            d = docx.Document(entry.abs_path)
            chunks = [p.text for p in d.paragraphs]
            for t in d.tables:
                for row in t.rows:
                    chunks.append(" ".join(c.text for c in row.cells))
            return "\n".join(chunks)
        if entry.ext == "pptx":
            prs = Presentation(entry.abs_path)
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            chunks.append(" ".join(c.text for c in row.cells))
            return "\n".join(chunks)
        if entry.ext in ("md", "txt", "csv", "tsv", "json", "py"):
            with open(entry.abs_path, encoding="utf-8", errors="ignore") as f:
                return f.read()
    except Exception:
        return ""
    return ""


def harvest_date_candidates(project_files: list[FileEntry]) -> list[str]:
    """YYYYMMDD strings found in filenames/content of non-encrypted sibling files."""
    candidates: set[str] = set()
    for f in project_files:
        if f.is_encrypted:
            continue
        text = f.name + "\n" + _extract_plain_text(f)
        for m in _DATE_RE.finditer(text):
            y, mo, da = m.groups()
            try:
                candidates.add(f"{int(y):04d}{int(mo):02d}{int(da):02d}")
            except ValueError:
                continue
    return sorted(candidates)


def generate_candidate_passwords(
    abbrev: str,
    date_candidates: list[str],
    ext: str,
    extra_ext_variants: list[str] | None = None,
) -> list[str]:
    """DA-{abbrev}-{date}-{ext} plus a handful of casing/format variants."""
    abbrev_variants = {abbrev, abbrev.upper(), abbrev.lower(), abbrev.capitalize()}
    ext_variants = {ext, ext.upper(), ext.lower(), f".{ext}"}
    if extra_ext_variants:
        ext_variants.update(extra_ext_variants)

    out = []
    for ab, date_str, ec in itertools.product(abbrev_variants, date_candidates, ext_variants):
        out.append(f"DA-{ab}-{date_str}-{ec}")
    return out


def _load_password_cache() -> dict:
    if _PASSWORD_CACHE_PATH.exists():
        return json.loads(_PASSWORD_CACHE_PATH.read_text(encoding="utf-8"))
    return {}


def _save_password_cache(cache: dict) -> None:
    _PASSWORD_CACHE_PATH.write_text(json.dumps(cache, ensure_ascii=False, indent=2), encoding="utf-8")


def try_passwords(entry: FileEntry, passwords: list[str]) -> tuple[bytes | None, str | None]:
    """Try each password in order; return (decrypted_bytes, password) on success."""
    for pw in passwords:
        try:
            with open(entry.abs_path, "rb") as fh:
                of = msoffcrypto.OfficeFile(fh)
                of.load_key(password=pw)
                dec = io.BytesIO()
                of.decrypt(dec)
            return dec.getvalue(), pw
        except Exception:
            continue
    return None, None


@dataclass
class DecryptResult:
    success: bool
    data: bytes | None
    password: str | None
    candidates_tried: int


def decrypt_entry(entry: FileEntry, catalog: Catalog, glossary: Glossary, extra_passwords: list[str] | None = None) -> DecryptResult:
    """Full pipeline: resolve project abbreviation, harvest dates from
    siblings, brute-force the documented password format, optionally try
    caller-supplied extra passwords too."""
    if not entry.is_encrypted:
        with open(entry.abs_path, "rb") as f:
            return DecryptResult(True, f.read(), None, 0)

    cache = _load_password_cache()
    cache_key = entry.project_key
    passwords: list[str] = []
    if cache_key in cache:
        passwords.append(cache[cache_key])

    project = next((p for p in catalog.projects if p.key == entry.project_key), None)
    alias = next((p for p in glossary.projects if p.full_name == entry.project_key), None)
    if project and alias:
        date_candidates = harvest_date_candidates(project.files)
        passwords.extend(generate_candidate_passwords(alias.code, date_candidates, entry.ext))

    if extra_passwords:
        passwords.extend(extra_passwords)

    data, pw = try_passwords(entry, passwords)
    if data is not None and pw is not None:
        cache[cache_key] = pw
        _save_password_cache(cache)
        return DecryptResult(True, data, pw, len(passwords))

    return DecryptResult(False, None, None, len(passwords))


def is_locked(path, ext: str) -> bool:
    return ext.lower().lstrip(".") in ("docx", "xlsx", "pptx") and sniff_is_ole(path)
