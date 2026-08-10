"""python-pptx based parser: slide shapes (including nested groups), tables,
speaker notes, and embedded pictures -- with run-level formatting.

python-pptx's high-level API doesn't expose character highlight color, so
we reach into the run's raw XML (`<a:highlight>`) for that. Font color is
best-effort: theme-referenced colors (as opposed to explicit RGB) are
reported as None rather than resolved, since resolving them requires
walking the theme part.
"""

from __future__ import annotations

import hashlib
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from datastel_rag.ingest.models import Block, FormattedRun, ImageAsset, ParsedDocument

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _run_highlight(run) -> str | None:
    rPr = run._r.find(f"{{{_A_NS}}}rPr")
    if rPr is None:
        return None
    hl = rPr.find(f"{{{_A_NS}}}highlight")
    if hl is None:
        return None
    srgb = hl.find(f"{{{_A_NS}}}srgbClr")
    if srgb is not None:
        return srgb.get("val")
    return None


def _run_format(run) -> FormattedRun:
    color = None
    try:
        if run.font.color and run.font.color.type is not None and run.font.color.rgb is not None:
            color = str(run.font.color.rgb)
    except Exception:
        pass
    return FormattedRun(
        text=run.text,
        bold=bool(run.font.bold),
        italic=bool(run.font.italic),
        underline=bool(run.font.underline),
        font_color=color,
        highlight=_run_highlight(run),
    )


def _text_frame_blocks(tf, block_id_prefix: str, location: dict) -> list[Block]:
    blocks = []
    for pi, para in enumerate(tf.paragraphs):
        text = "".join(r.text for r in para.runs) or para.text
        if not text.strip():
            continue
        runs = [_run_format(r) for r in para.runs if r.text]
        blocks.append(
            Block(
                block_id=f"{block_id_prefix}/para{pi}",
                kind="paragraph",
                text=text,
                runs=runs,
                location=location,
            )
        )
    return blocks


def _table_blocks(table, block_id_prefix: str, location: dict) -> list[Block]:
    blocks = []
    for ri, row in enumerate(table.rows):
        text = " | ".join(c.text for c in row.cells)
        if not text.strip():
            continue
        loc = dict(location)
        loc["row"] = ri
        blocks.append(Block(block_id=f"{block_id_prefix}/row{ri}", kind="table_row", text=text, location=loc))
    return blocks


def _save_image(image, cache_dir: Path, image_id: str) -> str:
    cache_dir.mkdir(parents=True, exist_ok=True)
    ext = image.ext or "png"
    digest = hashlib.sha1(image.blob).hexdigest()[:12]
    out_path = cache_dir / f"{image_id.replace('/', '_')}_{digest}.{ext}"
    if not out_path.exists():
        out_path.write_bytes(image.blob)
    return str(out_path)


def _walk_shapes(shapes, slide_idx: int, block_id_prefix: str, doc: ParsedDocument, image_cache_dir: Path):
    for si, shape in enumerate(shapes):
        sid = f"{block_id_prefix}/shape{si}"
        loc = {"slide": slide_idx, "shape_name": getattr(shape, "name", "")}

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk_shapes(shape.shapes, slide_idx, sid, doc, image_cache_dir)
            continue

        if shape.has_text_frame:
            doc.blocks.extend(_text_frame_blocks(shape.text_frame, sid, loc))

        if shape.has_table:
            doc.blocks.extend(_table_blocks(shape.table, sid, loc))

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                path = _save_image(shape.image, image_cache_dir, sid)
                doc.images.append(
                    ImageAsset(
                        image_id=sid,
                        location=loc,
                        cache_path=path,
                        width=Emu(shape.width).inches if shape.width else None,
                        height=Emu(shape.height).inches if shape.height else None,
                    )
                )
            except Exception as e:
                doc.parse_errors.append(f"image extract failed at {sid}: {e}")


def parse_pptx(abs_path: str, source_rel_path: str, project_key: str, image_cache_dir: Path) -> ParsedDocument:
    doc = ParsedDocument(source_rel_path=source_rel_path, project_key=project_key, ext="pptx")
    try:
        prs = Presentation(abs_path)
    except Exception as e:
        doc.parse_errors.append(f"open failed: {e}")
        return doc

    for slide_idx, slide in enumerate(prs.slides):
        _walk_shapes(slide.shapes, slide_idx, f"slide{slide_idx}", doc, image_cache_dir)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                doc.blocks.append(
                    Block(block_id=f"slide{slide_idx}/notes", kind="note", text=notes, location={"slide": slide_idx})
                )

    doc.meta = {"num_slides": len(prs.slides), "slide_width_emu": prs.slide_width, "slide_height_emu": prs.slide_height}
    return doc
